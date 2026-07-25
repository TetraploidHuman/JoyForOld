# -*- coding: utf-8 -*-
"""Generate JoyForOld competition PPT for judges."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — formal competition style
NAVY = RGBColor(0x0B, 0x2C, 0x4A)
BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT = RGBColor(0x1A, 0x9B, 0xB8)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x4A, 0x55, 0x68)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD0, 0xD7, 0xDE)
SOFT_BLUE = RGBColor(0xE8, 0xF3, 0xFB)


def set_run_font(run, size=18, bold=False, color=DARK, name="微软雅黑"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    # Also set East Asian font for Chinese
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        {"typeface": name},
    )
    # remove existing ea if any
    for child in list(rPr):
        if child.tag.endswith("}ea"):
            rPr.remove(child)
    rPr.append(ea)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color, name=font_name)
    return box


def add_paragraph(tf, text, size=16, bold=False, color=DARK, space_before=6,
                  space_after=4, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return p


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def slide_chrome(slide, title, page_no, total, subtitle=None):
    """Standard content slide header/footer."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BLUE)
    add_rect(slide, 0, Inches(1.05), SLIDE_W, Inches(0.015), LINE)
    add_textbox(slide, Inches(0.6), Inches(0.28), Inches(10.5), Inches(0.5),
                title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.72), Inches(11), Inches(0.3),
                    subtitle, size=12, color=GRAY)
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), LIGHT)
    add_textbox(slide, Inches(0.6), SLIDE_H - Inches(0.32), Inches(8), Inches(0.28),
                "JoyForOld  |  面向老年用户的端云协同 GUI 智能体系统",
                size=10, color=GRAY)
    add_textbox(slide, Inches(11.2), SLIDE_H - Inches(0.32), Inches(1.5), Inches(0.28),
                f"{page_no} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, title, body_lines, title_size=14, body_size=12):
    add_round_rect(slide, left, top, width, height, SOFT_BLUE)
    add_rect(slide, left, top, Inches(0.08), height, BLUE)
    box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, size=title_size, bold=True, color=NAVY)
    for line in body_lines:
        add_paragraph(tf, line, size=body_size, color=GRAY, space_before=4, space_after=2)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 14

    # ---------- 1 Cover ----------
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.85), SLIDE_W, Inches(1.65), BLUE)
    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
                "计算机软件开发类竞赛作品", size=16, color=RGBColor(0xA8, 0xC5, 0xDE))
    add_textbox(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.0),
                "JoyForOld", size=54, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.6),
                "面向老年用户的端云协同 GUI 智能体系统", size=24, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.5),
                "自然语言驱动 · 跨应用界面操控 · 安全确认 · 家庭远程协同",
                size=14, color=RGBColor(0xC5, 0xD8, 0xE8))
    add_textbox(s, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.35),
                "Android 智能体软件  |  人工智能 × 人机交互 × 智慧助老",
                size=14, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.35),
                "版本 1.0  ·  包名 com.tetraploid.joyforold",
                size=12, color=RGBColor(0xD0, 0xE8, 0xF8))

    # ---------- 2 Agenda ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "目录 CONTENTS", 2, total, "评审汇报结构")
    items = [
        ("01", "研究背景与问题提出"),
        ("02", "作品定位与能力闭环"),
        ("03", "功能体系与典型场景"),
        ("04", "系统架构与技术路线"),
        ("05", "核心技术创新点"),
        ("06", "方案对比与社会价值"),
        ("07", "总结与展望"),
    ]
    for i, (num, title) in enumerate(items):
        y = Inches(1.35) + Inches(0.72) * i
        add_round_rect(s, Inches(1.2), y, Inches(10.8), Inches(0.58), SOFT_BLUE if i % 2 == 0 else LIGHT)
        add_textbox(s, Inches(1.5), y + Inches(0.1), Inches(1.2), Inches(0.4),
                    num, size=20, bold=True, color=BLUE)
        add_textbox(s, Inches(3.0), y + Inches(0.12), Inches(8.5), Inches(0.4),
                    title, size=18, color=NAVY)

    # ---------- 3 Background ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "01  研究背景", 3, total, "老龄化 × 智能手机普及 × 数字鸿沟")
    card(s, Inches(0.5), Inches(1.35), Inches(3.9), Inches(2.35),
         "社会背景",
         ["人口老龄化持续加深", "家庭结构小型化、子女异地就业", "养老与照护日益依赖数字化手段"],
         title_size=16, body_size=13)
    card(s, Inches(4.7), Inches(1.35), Inches(3.9), Inches(2.35),
         "现实矛盾",
         ["政务、医疗、出行、通讯、支付", "加速向移动端迁移", "「有设备」≠「会使用」"],
         title_size=16, body_size=13)
    card(s, Inches(8.9), Inches(1.35), Inches(3.9), Inches(2.35),
         "家庭现状",
         ["子女视频口授成本高", "界面状态难共享、难复用", "紧急场景更容易失败"],
         title_size=16, body_size=13)

    box = s.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "关键洞察"
    set_run_font(run, size=16, bold=True, color=BLUE)
    points = [
        "老年用户面临界面复杂、路径长、易中断、改版后失效等问题，构成典型数字鸿沟。",
        "主流语音助手擅长问答与系统技能，难以完成跨应用、多步骤的真实界面操作。",
        "智慧助老需要从「陪聊式陪伴」走向「任务完成型智能辅助」。",
        "大模型规划能力 + 无障碍自动化 + 端侧轻量推理，为银发 GUI 智能体提供技术条件。",
    ]
    for t in points:
        add_paragraph(tf, "●  " + t, size=14, color=DARK, space_before=8, space_after=2)

    # ---------- 4 Problem ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "01  核心问题定义", 4, total, "本作品要回答的科学/工程问题")
    add_round_rect(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.55), SOFT_BLUE)
    add_textbox(s, Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.3),
                "在保障安全与可控的前提下，如何将老年用户的自然语言意图，稳定转化为对真实智能手机界面"
                "（含第三方应用）的多步骤可执行操作，并在个人能力不足时提供家庭远程托底？",
                size=16, bold=True, color=NAVY)

    subs = [
        ("感知", "如何可靠获取界面状态\n（控件树 / 视觉）"),
        ("理解", "如何理解口语化、\n省略式老年指令"),
        ("决策", "如何生成动作序列\n并支持失败重规划"),
        ("安全", "如何防止误发、\n误拨等高风险操作"),
        ("协同", "如何引入家属远程\n协助而不增负担"),
    ]
    for i, (t, b) in enumerate(subs):
        x = Inches(0.5) + Inches(2.5) * i
        add_round_rect(s, x, Inches(3.3), Inches(2.35), Inches(2.7), LIGHT)
        add_rect(s, x, Inches(3.3), Inches(2.35), Inches(0.5), BLUE)
        add_textbox(s, x, Inches(3.38), Inches(2.35), Inches(0.4),
                    t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.12), Inches(4.0), Inches(2.1), Inches(1.8),
                    b, size=13, color=DARK, align=PP_ALIGN.CENTER)

    # ---------- 5 Positioning ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "02  作品定位", 5, total, "端云协同 GUI 智能体，而非陪聊机器人")
    add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
                "JoyForOld：面向老年用户的端云协同 GUI Agent 系统",
                size=20, bold=True, color=NAVY)
    add_textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.7),
                "以自然语言为统一入口，以无障碍界面感知与自动化为执行载体，以端云分层决策为智能核心，"
                "以语音二次确认与动作守卫为安全约束，以家属远程协助为照护托底。",
                size=14, color=GRAY)

    goals = [
        ("可用性", "一句话办事\n降低操作负担"),
        ("可靠性", "高频场景可复现\n减少半途失败"),
        ("安全性", "敏感操作确认\n抑制误触风险"),
        ("可达性", "弱网保留快路径\n端侧意图执行"),
        ("协同性", "老人自主使用\n家属远程托底"),
        ("可扩展", "模型可配置\n任务可扩展"),
    ]
    for i, (t, b) in enumerate(goals):
        col, row = i % 3, i // 3
        x = Inches(0.5) + Inches(4.2) * col
        y = Inches(2.7) + Inches(1.85) * row
        add_round_rect(s, x, y, Inches(4.0), Inches(1.65), SOFT_BLUE)
        add_textbox(s, x + Inches(0.25), y + Inches(0.25), Inches(3.5), Inches(0.4),
                    t, size=18, bold=True, color=BLUE)
        add_textbox(s, x + Inches(0.25), y + Inches(0.7), Inches(3.5), Inches(0.8),
                    b, size=14, color=DARK)

    # ---------- 6 Capability loop ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "02  能力闭环", 6, total, "听懂 → 看懂 → 想清 → 做对 → 问准 / 帮托")
    steps = [
        ("听懂", "唤醒 + ASR\n本地唤醒词 / 流式识别"),
        ("看懂", "界面观测\n控件树优先 / 视觉回退"),
        ("想清", "分层决策\n模板·NLU·Agent"),
        ("做对", "自动执行\n点击 / 输入 / 滚动"),
        ("问准", "安全确认\n发送·拨号二次确认"),
        ("帮托", "家庭协同\n协助码远程指令"),
    ]
    for i, (t, b) in enumerate(steps):
        x = Inches(0.35) + Inches(2.15) * i
        add_round_rect(s, x, Inches(2.0), Inches(2.0), Inches(3.6), LIGHT)
        add_rect(s, x, Inches(2.0), Inches(2.0), Inches(0.7), BLUE if i < 4 else ACCENT)
        add_textbox(s, x, Inches(2.15), Inches(2.0), Inches(0.5),
                    t, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(3.0), Inches(1.8), Inches(2.3),
                    b, size=13, color=DARK, align=PP_ALIGN.CENTER)
        if i < 5:
            add_textbox(s, x + Inches(1.85), Inches(3.4), Inches(0.35), Inches(0.4),
                        "→", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
                "评价标准：事情是否办成，而非对话是否流畅。任务闭环优先于生成炫技。",
                size=14, bold=True, color=NAVY)

    # ---------- 7 Functions ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "03  功能体系", 7, total, "覆盖入口、理解、执行、安全与协同")
    funcs = [
        ("多通道入口", "唤醒词 / 语音按钮\n文本输入 / 建议指令"),
        ("跨应用操控", "无障碍读屏执行\n微信支持组件协同"),
        ("端侧快路径", "模板 + 离线 NLU\n系统 Intent 直达"),
        ("云端 Agent", "工具调用式规划\n观察—执行—再规划"),
        ("安全确认", "发送 / 拨号确认\n动作守卫防死循环"),
        ("悬浮 + IME", "跨应用悬浮助手\n专用输入提升成功率"),
        ("家庭协助", "协助码配对\n看屏 / 远端指令"),
        ("银发刚需包", "通话·导航·读消息\n健康码·紧急求助"),
    ]
    for i, (t, b) in enumerate(funcs):
        col, row = i % 4, i // 4
        x = Inches(0.45) + Inches(3.2) * col
        y = Inches(1.4) + Inches(2.55) * row
        add_round_rect(s, x, y, Inches(3.05), Inches(2.3), SOFT_BLUE)
        add_textbox(s, x + Inches(0.2), y + Inches(0.35), Inches(2.65), Inches(0.5),
                    t, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), y + Inches(1.0), Inches(2.65), Inches(1.0),
                    b, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # ---------- 8 Scenarios ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "03  典型应用场景", 8, total, "可现场演示的刚需任务")
    scenes = [
        ("01", "亲情通话", "「打电话给女儿」→ 联系人解析 → 拨号前确认 → 执行拨号"),
        ("02", "微信代发", "自然语言描述收件人与内容 → 自动操作 → 发送前强制确认"),
        ("03", "导航回家", "「我要回家」→ 目的地解析 → 调起地图导航深链"),
        ("04", "弱网设置", "「打开 WiFi / 放大字体」→ 端侧意图命中 → 本地快速完成"),
        ("05", "信息获取", "读未读消息、询问时间等 → 视觉信息转化为听觉反馈"),
        ("06", "远程托底", "协助码配对 → 家属看屏 / 发指令 → 老人手机落地执行"),
    ]
    for i, (n, t, b) in enumerate(scenes):
        y = Inches(1.3) + Inches(0.85) * i
        add_rect(s, Inches(0.55), y, Inches(0.7), Inches(0.7), BLUE)
        add_textbox(s, Inches(0.55), y + Inches(0.15), Inches(0.7), Inches(0.45),
                    n, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.45), y + Inches(0.05), Inches(2.2), Inches(0.55),
                    t, size=16, bold=True, color=NAVY)
        add_textbox(s, Inches(3.7), y + Inches(0.08), Inches(9.0), Inches(0.55),
                    b, size=13, color=DARK)

    # ---------- 9 Architecture ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "04  系统架构", 9, total, "分层解耦 · 端云协同 · 执行可观测")
    layers = [
        ("交互呈现层", "Compose 主界面 / 悬浮窗 / 通知确认 / TTS 播报"),
        ("运行时聚合层", "AgentRuntime：会话、权限、语音与协助状态统一调度"),
        ("决策路由层", "模板 → 离线 NLU → 系统 Intent → 预设 → 云端 Agent"),
        ("执行与感知层", "AccessibilityGateway / 无障碍服务 / 输入法 / 系统深链"),
        ("协同与模型层", "协助协议与中继服务  ·  唤醒 / VAD / 意图 ONNX 资源"),
    ]
    for i, (t, b) in enumerate(layers):
        y = Inches(1.35) + Inches(0.95) * i
        add_round_rect(s, Inches(1.0), y, Inches(11.3), Inches(0.82), SOFT_BLUE if i % 2 == 0 else LIGHT)
        add_rect(s, Inches(1.0), y, Inches(2.6), Inches(0.82), BLUE)
        add_textbox(s, Inches(1.05), y + Inches(0.2), Inches(2.5), Inches(0.5),
                    t, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(3.85), y + Inches(0.22), Inches(8.2), Inches(0.5),
                    b, size=14, color=DARK)

    # ---------- 10 Tech ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "04  技术路线", 10, total, "关键技术选型与处理流水线")
    techs = [
        ("客户端", "Kotlin · Compose\nKoin · Coroutines"),
        ("界面操控", "Accessibility\n跨应用读屏执行"),
        ("端侧语音", "Sherpa KWS\nSilero VAD"),
        ("端侧 NLU", "RoBERTa-Mini\nINT8 ONNX"),
        ("云端能力", "豆包流式 ASR\n可配置 LLM"),
        ("家庭协同", "Ktor · JWT\nWebSocket 中继"),
    ]
    for i, (t, b) in enumerate(techs):
        x = Inches(0.45) + Inches(2.15) * i
        add_round_rect(s, x, Inches(1.35), Inches(2.05), Inches(2.2), SOFT_BLUE)
        add_textbox(s, x + Inches(0.1), Inches(1.55), Inches(1.85), Inches(0.4),
                    t, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(2.1), Inches(1.85), Inches(1.2),
                    b, size=12, color=DARK, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.4),
                "主处理流水线", size=16, bold=True, color=NAVY)
    pipeline = (
        "麦克风采集 → VAD 门控 → 唤醒检出 → 预缓冲衔接 ASR → 多层路由 → "
        "本地执行 / Agent 循环 → 安全确认 → 无障碍执行 → TTS 反馈"
    )
    add_round_rect(s, Inches(0.5), Inches(4.35), Inches(12.3), Inches(1.5), LIGHT)
    add_textbox(s, Inches(0.75), Inches(4.65), Inches(11.8), Inches(1.0),
                pipeline, size=15, color=DARK)

    # ---------- 11 Innovation ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "05  核心创新点（一）", 11, total, "范式 · 架构 · 端侧智能 · 安全")
    inns = [
        ("01  GUI 智能体范式", "将大模型从对话生成推进到界面操控闭环，以事情办成为评价标准。"),
        ("02  端云多层路由", "模板 / 离线 NLU / 系统意图 / 预设 / 云端 Agent 级联分流。"),
        ("03  端侧慎行 NLU", "INT8 意图分类 + 置信度门控：自动执行 / 澄清 / 回退。"),
        ("04  唤醒链路耦合", "KWS + VAD + 预缓冲 ASR，缩短唤起至表达的交互断裂。"),
        ("05  语音二次确认", "发送、拨号等高风险动作强制口语确认，安全优先于全自动。"),
        ("06  动作守卫机制", "抑制重复无效点击与死循环，保持代操作的可控性。"),
    ]
    for i, (t, b) in enumerate(inns):
        col, row = i % 2, i // 2
        x = Inches(0.5) + Inches(6.4) * col
        y = Inches(1.3) + Inches(1.7) * row
        add_round_rect(s, x, y, Inches(6.15), Inches(1.5), SOFT_BLUE)
        add_textbox(s, x + Inches(0.25), y + Inches(0.2), Inches(5.7), Inches(0.4),
                    t, size=15, bold=True, color=NAVY)
        add_textbox(s, x + Inches(0.25), y + Inches(0.7), Inches(5.7), Inches(0.65),
                    b, size=13, color=DARK)

    # ---------- 12 Innovation 2 ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "05  核心创新点（二）", 12, total, "感知 · 工程可靠 · 照护协同 · 隐私")
    inns2 = [
        ("07  双通道界面感知", "控件树精确点击为主，截图多模态规划为辅，适配复杂第三方界面。"),
        ("08  Action-set 工程化", "微信 / 地图等高频流程剧本化，智能规划与稳定路径相结合。"),
        ("09  家庭照护闭环", "协助协议 + 中继服务，形成老人自主与家属托底的三层能力结构。"),
        ("10  隐私脱敏上云", "页面上下文脱敏、密钥本地保存，平衡能力上云与隐私保护。"),
    ]
    for i, (t, b) in enumerate(inns2):
        y = Inches(1.35) + Inches(1.2) * i
        add_round_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.05), SOFT_BLUE)
        add_rect(s, Inches(0.6), y, Inches(0.12), Inches(1.05), BLUE)
        add_textbox(s, Inches(1.0), y + Inches(0.15), Inches(11.4), Inches(0.35),
                    t, size=16, bold=True, color=NAVY)
        add_textbox(s, Inches(1.0), y + Inches(0.52), Inches(11.4), Inches(0.4),
                    b, size=14, color=DARK)

    # ---------- 13 Comparison + Value ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "06  方案对比与社会价值", 13, total, "差异化竞争力与民生导向")
    add_textbox(s, Inches(0.55), Inches(1.25), Inches(12), Inches(0.35),
                "与现有方案相比：本作品强在「系统集成 + 适老化约束 + 家庭协同」",
                size=14, bold=True, color=BLUE)

    # Simple comparison rows
    headers = ["维度", "语音助手/陪聊", "脚本 RPA", "本作品"]
    rows = [
        ["跨应用深度操作", "有限", "强（脚本内）", "智能体 + 动作集"],
        ["适老化设计", "弱—中", "弱", "强"],
        ["安全确认", "一般", "弱", "语音二次确认"],
        ["弱网快路径", "弱", "强", "端侧 NLU/模板"],
        ["家庭远程托底", "少见", "少见", "协助闭环"],
    ]
    # header
    y0 = Inches(1.7)
    widths = [Inches(2.4), Inches(3.0), Inches(3.0), Inches(3.5)]
    xs = [Inches(0.7)]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)
    for i, h in enumerate(headers):
        add_rect(s, xs[i], y0, widths[i], Inches(0.45), BLUE)
        add_textbox(s, xs[i], y0 + Inches(0.08), widths[i], Inches(0.35),
                    h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + Inches(0.45) * (r + 1)
        bg = LIGHT if r % 2 == 0 else WHITE
        for i, cell in enumerate(row):
            add_rect(s, xs[i], y, widths[i], Inches(0.45), bg)
            # draw light border effect via darker text only
            col = NAVY if i == 0 or i == 3 else GRAY
            bold = i == 0 or i == 3
            add_textbox(s, xs[i], y + Inches(0.08), widths[i], Inches(0.35),
                        cell, size=12, bold=bold, color=col, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.35),
                "社会价值", size=15, bold=True, color=NAVY)
    vals = [
        "缓解老年数字鸿沟，提升独立完成生活事务的能力；",
        "减轻子女视频口授负担，形成可复用的家庭远程协助机制；",
        "提供无需逐一改造业务 App 的第三方适老化能力层思路；",
        "可作为 GUI Agent、老年 HCI 与智慧养老方向的可演示平台。",
    ]
    box = s.shapes.add_textbox(Inches(0.7), Inches(4.95), Inches(12), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for v in vals:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = "●  " + v
        set_run_font(run, size=13, color=DARK)

    # ---------- 14 Summary ----------
    s = prs.slides.add_slide(blank)
    slide_chrome(s, "07  总结", 14, total, "问题导向 · 技术贯通 · 工程可演示")
    add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
                "JoyForOld 以适老化任务闭环为导向，完成从语言理解到跨应用界面执行的系统贯通，"
                "并以端云分流与安全门控保证可用性与可控性。",
                size=15, color=DARK)

    three = [
        ("问题硬", "数字鸿沟与家庭口授低效\n是真实刚需，而非虚构场景"),
        ("技术新", "GUI Agent + 端云路由\n+ 安全确认 + 家庭协同"),
        ("落地实", "真机可演示、模块可扩展\n边界清晰、工程完备"),
    ]
    for i, (t, b) in enumerate(three):
        x = Inches(0.55) + Inches(4.2) * i
        add_round_rect(s, x, Inches(2.4), Inches(3.95), Inches(2.4), SOFT_BLUE)
        add_rect(s, x, Inches(2.4), Inches(3.95), Inches(0.65), BLUE)
        add_textbox(s, x, Inches(2.52), Inches(3.95), Inches(0.45),
                    t, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), Inches(3.3), Inches(3.55), Inches(1.3),
                    b, size=14, color=DARK, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
                "听得懂 · 看得见 · 点得到 · 问得准 · 家人能帮",
                size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.4),
                "感谢各位评委老师批评指正",
                size=16, color=NAVY, align=PP_ALIGN.CENTER)

    out = Path(__file__).resolve().parent / "JoyForOld-Judge-PPT.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    build()
